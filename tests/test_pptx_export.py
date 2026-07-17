import os
import sys

import pytest

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

import config
from core import finance, pptx_export


def test_template_disponible_reflete_presence_fichier(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(config, "TEMPLATE_PPTX_PATH", tmp_path / "absent.pptx")
    assert pptx_export.template_disponible() is False

    fichier = tmp_path / "present.pptx"
    fichier.write_bytes(b"")
    monkeypatch.setattr(config, "TEMPLATE_PPTX_PATH", fichier)
    assert pptx_export.template_disponible() is True


def test_generer_pptx_sans_template_leve_erreur_claire(tmp_path, monkeypatch) -> None:
    monkeypatch.setattr(config, "TEMPLATE_PPTX_PATH", tmp_path / "absent.pptx")
    with pytest.raises(FileNotFoundError, match="Template PowerPoint introuvable"):
        pptx_export.generer_pptx(
            demande={"titre": "Mission", "reference": "RFX001", "client_nom": "C"},
            lignes=[],
            chemin_sortie=tmp_path / "sortie.pptx",
        )


def _texte_slides(prs) -> str:
    return " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in pptx_export._formes(slide)
        if shape.has_text_frame
    )


def _texte_tables(prs) -> str:
    return " ".join(
        cell.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_table
        for row in shape.table.rows
        for cell in row.cells
    )


@pytest.mark.skipif(
    not config.TEMPLATE_PPTX_PATH.exists(),
    reason="data/template_proposition.pptx est un fichier confidentiel local, absent en CI",
)
def test_generer_pptx_avec_le_vrai_template(tmp_path) -> None:
    from pptx import Presentation

    demande = {"titre": "Outil de démonstration RFX000TEST", "reference": "RFX000TEST",
               "client_nom": "ClientDemo"}
    lignes = [
        {"prenom": "Alice", "nom": "Dupont", "grade": "Manager 2 (M2)", "seniorite": None,
         "titre": "Actuaire Senior", "nb_jours": 20, "tjm_applique": 1370,
         "annees_experience": 8, "formation": "Master Actuariat, ISFA Lyon",
         "synthese_cv": "Synthèse ciblée mission pour Alice.", "photo_path": None,
         "cv_complet_json": '{"competences": ["Solvabilité 2"], "experiences": []}'},
        {"prenom": "Bob", "nom": "Martin", "grade": "Junior 1 (J1)", "seniorite": None,
         "titre": None, "nb_jours": 15, "tjm_applique": 820,
         "annees_experience": None, "formation": None, "synthese_cv": None, "photo_path": None,
         "cv_complet_json": '{"experiences": [{"role": "Analyste", "client": "MutuelleX", "description": "Reporting QRT"}]}'},
    ]
    contenu = {
        "contexte_redige": "Le client souhaite fiabiliser son processus RFX000TEST.",
        "demarche": {
            "cadrage": "Ateliers de cadrage de démonstration.",
            "analyse": "Analyse de démonstration.",
            "realisation": "Réalisation de démonstration.",
            "accompagnement": "Accompagnement de démonstration.",
            "restitution": "Restitution de démonstration.",
        },
    }
    redacteur = {"nom": "Jean Dupuis", "fonction": "Manager",
                 "email": "jean.dupuis@actuelia.fr", "telephone": "06 00 00 00 00"}
    sortie = tmp_path / "generation.pptx"

    total = pptx_export.generer_pptx(demande=demande, lignes=lignes, chemin_sortie=sortie,
                                     contenu=contenu, redacteur=redacteur)

    assert total == finance.total_mission(lignes)
    assert sortie.exists()

    prs = Presentation(sortie)
    textes = _texte_slides(prs)
    tables = _texte_tables(prs)

    # Page de garde : référence, intitulé, date, rédacteur (marqueurs remplacés).
    garde = " ".join(s.text_frame.text for s in pptx_export._formes(prs.slides[0]) if s.has_text_frame)
    assert "RFX000TEST" in garde
    assert "Outil de démonstration RFX000TEST" in garde
    assert "Jean Dupuis" in garde
    assert "[Prénom NOM]" not in garde
    # Contexte rédigé.
    assert contenu["contexte_redige"] in textes
    # Démarche : slide dédiée (titre + 5 phases en texte), la frise est retirée.
    assert "Cadrage :" in textes
    assert "Restitution :" in textes
    assert "Restitution de démonstration." in textes
    assert "[Phase 1]" not in textes  # frise retirée
    slides_demarche = [
        s for s in prs.slides
        if any(sh.has_text_frame and sh.text_frame.text.strip() == "Démarche opérationnelle"
               for sh in pptx_export._formes(s))
    ]
    assert len(slides_demarche) == 1  # une seule slide démarche dédiée
    txt_demarche = " ".join(
        sh.text_frame.text for sh in pptx_export._formes(slides_demarche[0]) if sh.has_text_frame
    )
    assert "Cadrage :" in txt_demarche
    assert "Restitution de démonstration." in txt_demarche
    # Client injecté dans les modalités.
    assert "ClientDemo" in textes
    # Fiches CV : une par consultant, avec leur vrai nom (marqueur remplacé sur ces slides).
    assert "Alice DUPONT" in textes
    assert "Bob MARTIN" in textes
    assert "Synthèse ciblée mission pour Alice." in tables
    assert "Reporting QRT" in tables
    # Les images du collaborateur-exemple du template sont retirées des fiches CV.
    for slide_cv in [s for s in prs.slides if s.slide_layout.name == "Slide CV"]:
        assert [s for s in slide_cv.shapes if s.shape_type == 13] == []
    # Budget : une ligne par consultant + totaux dans l'intro.
    assert "Manager 2 (M2)" in tables
    assert "1 370 €" in tables
    assert "35 jours" in textes  # 20 + 15
    assert "39 700" in textes    # total mission

    # Équipe projet (slide Modalités) : un bloc par consultant, superviseur retiré,
    # aucun marqueur de nom résiduel.
    modalites = " ".join(
        s.text_frame.text for s in pptx_export._formes(prs.slides[6]) if s.has_text_frame
    )
    assert "Alice DUPONT" in modalites
    assert "Bob MARTIN" in modalites
    assert "Solvabilité 2" in modalites  # expertise = compétences
    # Valeur ajoutée = début de la synthèse ciblée mission (bornée à 2 phrases).
    assert "Valeur ajoutée" in modalites
    assert "Synthèse ciblée mission pour Alice." in modalites
    assert "Superviseur" not in modalites
    assert "[Prénom NOM]" not in modalites
    # La démarche a migré sur sa propre slide : ni le pied ni les phases ne restent ici.
    assert "Démarche opérationnelle proposée" not in modalites
    assert "Restitution de démonstration." not in modalites

    # Régression : les parts « slide de notes » et « tags » appartiennent à UNE
    # seule slide, et tout r:id d'une slide doit résoudre dans ses relations.
    # Les enfreindre rend le fichier « endommagé » pour PowerPoint (pages
    # blanches, slides remplacées) — invisible pour python-pptx.
    import posixpath
    import zipfile
    from lxml import etree

    _R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    _EXCLUSIVES = ("/notesSlide", "/tags")
    proprietaires: dict[str, list[str]] = {}
    with zipfile.ZipFile(sortie) as z:
        for nom in z.namelist():
            if not (nom.startswith("ppt/slides/") and nom.endswith(".xml")):
                continue
            slide = posixpath.basename(nom)
            rels_nom = f"ppt/slides/_rels/{slide}.rels"
            rids = set()
            if rels_nom in z.namelist():
                for rel in etree.fromstring(z.read(rels_nom)):
                    rids.add(rel.get("Id"))
                    if rel.get("Type").endswith(_EXCLUSIVES):
                        cible = posixpath.normpath(posixpath.join("ppt/slides", rel.get("Target")))
                        proprietaires.setdefault(cible, []).append(slide)
            arbre = etree.fromstring(z.read(nom))
            pendantes = [
                v for el in arbre.iter() for a, v in el.attrib.items()
                if a.startswith(_R_NS) and v not in rids
            ]
            assert pendantes == [], f"{slide} : références r:id sans relation : {pendantes}"
    partagees = {part: slides for part, slides in proprietaires.items() if len(slides) > 1}
    assert partagees == {}, f"parts exclusives partagées entre slides : {partagees}"


@pytest.mark.skipif(
    not config.TEMPLATE_PPTX_PATH.exists(),
    reason="data/template_proposition.pptx est un fichier confidentiel local, absent en CI",
)
def test_generer_pptx_sans_contenu_ni_redacteur_conserve_les_marqueurs(tmp_path) -> None:
    from pptx import Presentation

    demande = {"titre": "Mission sans contenu", "reference": "RFX000VIDE", "client_nom": "C"}
    sortie = tmp_path / "generation_vide.pptx"

    pptx_export.generer_pptx(demande=demande, lignes=[], chemin_sortie=sortie,
                             contenu=None, redacteur=None)

    textes = _texte_slides(Presentation(sortie))
    # Sans contenu rédigé, la démarche garde ses marqueurs de phase du template.
    assert "[Phase 1]" in textes
    # Sans rédacteur, le bloc contact garde ses marqueurs à compléter.
    assert "[Prénom NOM]" in textes


@pytest.mark.skipif(
    not config.TEMPLATE_PPTX_PATH.exists(),
    reason="data/template_proposition.pptx est un fichier confidentiel local, absent en CI",
)
def test_generer_pptx_avec_cv_json_corrompu_ne_plante_pas(tmp_path) -> None:
    from pptx import Presentation

    lignes = [{"prenom": "Zoé", "nom": "Faure", "grade": None, "seniorite": "Senior",
               "titre": None, "nb_jours": 5, "tjm_applique": 1000,
               "annees_experience": None, "formation": None, "synthese_cv": None,
               "photo_path": None, "cv_complet_json": "{pas du json"}]
    sortie = tmp_path / "cv_corrompu.pptx"

    pptx_export.generer_pptx(demande={"titre": "T", "reference": "R", "client_nom": "C"},
                             lignes=lignes, chemin_sortie=sortie, contenu=None)

    assert "Zoé FAURE" in _texte_slides(Presentation(sortie))
