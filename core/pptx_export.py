"""Génération de la proposition PowerPoint à partir du template Actuelia.

Le fichier de template (data/template_proposition.pptx) contient la charte et
la trame de marque : il est volontairement hors Git (voir .gitignore et
README). Sans lui, la génération est indisponible.

Ce template est une trame « à trous » : ses 10 slides SONT la proposition, avec
des marqueurs [entre crochets] à remplir. On remplit donc les slides en place
(seule la fiche CV est dupliquée, une par consultant retenu) — au lieu de
reconstruire des slides. Le calcul financier reste 100% déterministe
(core/finance.py) : ce module ne met en forme que du texte et des chiffres déjà
calculés, jamais de calcul lui-même.
"""
import copy
import io
import json
import re
import zipfile
from datetime import date
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

import config
from . import finance
from .redaction import DEMARCHE_LABELS, PHASES_DEMARCHE

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Index (base 0) des slides de la trame.
_IDX_COUVERTURE = 0
_IDX_CONTEXTE = 5
_IDX_MODALITES = 6
_IDX_CV = 7
_IDX_BUDGET = 8


def template_disponible() -> bool:
    return config.TEMPLATE_PPTX_PATH.exists()


# --------------------------------------------------------------------------- #
#  Helpers bas niveau
# --------------------------------------------------------------------------- #
def _formes(conteneur):
    """Itère récursivement toutes les formes, groupes compris."""
    for shape in conteneur.shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from _formes(shape)


def _forme(slide, nom: str, *, recursif: bool = False):
    for shape in (_formes(slide) if recursif else slide.shapes):
        if shape.name == nom:
            return shape
    return None


def _formes_nommees(slide, nom: str) -> list:
    """Toutes les formes portant ce nom (le template réutilise « Tableau 3 »)."""
    return [s for s in slide.shapes if s.name == nom]


def _placeholder(slide, idx: int):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def _definir_texte(text_frame, lignes) -> None:
    """Remplace tout le texte en conservant la mise en forme du 1er run."""
    if isinstance(lignes, str):
        lignes = lignes.split("\n")
    lignes = [l for l in (lignes or [""])]

    for para in list(text_frame.paragraphs[1:]):
        para._p.getparent().remove(para._p)

    def _ecrire(paragraphe, texte: str) -> None:
        if paragraphe.runs:
            paragraphe.runs[0].text = texte
            for run in paragraphe.runs[1:]:
                run._r.getparent().remove(run._r)
        else:
            paragraphe.text = texte

    premier = text_frame.paragraphs[0]
    _ecrire(premier, lignes[0])
    corps = premier._p.getparent()
    for ligne in lignes[1:]:
        corps.append(copy.deepcopy(premier._p))
        _ecrire(text_frame.paragraphs[-1], ligne)


def _reduire_si_deborde(text_frame) -> None:
    """« Réduire le texte en cas de débordement » : PowerPoint rétrécit la police
    pour que le contenu (souvent long, généré par le LLM) tienne dans la forme."""
    try:
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _police_adaptee(nb_caracteres: int, nb_paragraphes: int, forme):
    """Taille de police (pt) pour que le texte tienne dans la forme, ou None si
    la taille du template (14 pt) suffit.

    L'auto-ajustement PowerPoint (_reduire_si_deborde) n'est recalculé qu'à
    l'ouverture du fichier dans PowerPoint et est ignoré par Google Slides : on
    fixe donc une taille explicite, estimée de façon déterministe (Calibri
    ≈ 0,5 × corps par caractère, interligne et espacements majorés)."""
    largeur_in = (forme.width or 0) / 914400
    hauteur_in = (forme.height or 0) / 914400
    if not largeur_in or not hauteur_in or not nb_caracteres:
        return None
    for pt in (14, 13, 12, 11, 10, 9, 8):
        caracteres_par_ligne = largeur_in * 144 / pt
        lignes = nb_caracteres / caracteres_par_ligne + nb_paragraphes
        hauteur = lignes * pt * 1.5 / 72 + nb_paragraphes * 12 / 72
        if hauteur <= hauteur_in * 0.95:
            return None if pt == 14 else Pt(pt)
    return Pt(8)


def _appliquer_police(text_frame, taille) -> None:
    if taille is None:
        return
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.size = taille


def _ajuster_police_zone(forme) -> None:
    """Fixe une taille de police explicite pour que le contenu de la zone tienne,
    puis laisse l'auto-ajustement PowerPoint affiner à l'ouverture."""
    tf = forme.text_frame
    texte = tf.text
    _appliquer_police(tf, _police_adaptee(len(texte), max(texte.count("\n") + 1, 1), forme))
    _reduire_si_deborde(tf)


def _reduire_police_cellule(text_frame, longueur: int) -> None:
    """Réduit la police d'une cellule selon la longueur du texte.

    Les cellules de tableau n'honorent pas l'auto-ajustement de PowerPoint
    (MSO_AUTO_SIZE) : on rétrécit donc explicitement pour les synthèses longues,
    par paliers déterministes, afin qu'elles tiennent dans la case CV.
    """
    if longueur > 900:
        taille = Pt(8)
    elif longueur > 600:
        taille = Pt(9)
    elif longueur > 350:
        taille = Pt(10)
    else:
        return
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.size = taille


def _remplacer_marqueurs(text_frame, remplacements: dict) -> None:
    """Remplace des marqueurs [xxx] au fil du texte, en conservant la mise en forme.

    Travaille au niveau du paragraphe (concaténation des runs) pour gérer les
    marqueurs que PowerPoint a éclatés sur plusieurs runs ; ne réécrit un
    paragraphe que s'il contient effectivement un marqueur.
    """
    for para in text_frame.paragraphs:
        plein = "".join(run.text for run in para.runs)
        if not plein:
            continue
        nouveau = plein
        for marqueur, valeur in remplacements.items():
            if marqueur in nouveau:
                nouveau = nouveau.replace(marqueur, valeur or "")
        if nouveau != plein:
            if para.runs:
                para.runs[0].text = nouveau
                for run in para.runs[1:]:
                    run._r.getparent().remove(run._r)
            else:
                para.text = nouveau


def _slide_vierge(prs, layout):
    """Nouvelle slide sur ce layout, débarrassée des placeholders auto-clonés."""
    dest = prs.slides.add_slide(layout)
    for shape in list(dest.shapes):
        shape._element.getparent().remove(shape._element)
    return dest


def _purger_refs_orphelines(element, rid_map: dict) -> None:
    """Retire de l'élément copié tout ce qui référence une relation non recopiée.

    Un r:id qui ne résout pas dans les relations de la slide de destination
    (tags d'add-in, notes…) suffit à faire passer le fichier pour endommagé
    aux yeux de PowerPoint."""
    for el in list(element.iter()):
        for attr, val in list(el.attrib.items()):
            if not attr.startswith("{%s}" % _R_NS):
                continue
            if val in rid_map:
                el.attrib[attr] = rid_map[val]
            else:
                parent = el.getparent()
                if parent is not None:
                    parent.remove(el)
                break
    # Une liste de tags vidée n'a plus de raison d'être.
    for el in list(element.iter(qn("p:custDataLst"))):
        if len(el) == 0:
            el.getparent().remove(el)


def _dupliquer_slide(prs, index_source):
    """Duplique prs.slides[index_source] en fin de présentation (images comprises).

    Seules les relations partageables (images, layout) sont recopiées : les
    parts qui appartiennent à UNE seule slide — slide de notes, tags d'add-in —
    ne doivent jamais être référencées par deux slides, sous peine de fichier
    « endommagé » pour PowerPoint. Les références à ces parts sont purgées du
    XML copié."""
    source = prs.slides[index_source]
    dest = _slide_vierge(prs, source.slide_layout)

    rid_map = {}
    for rid, rel in source.part.rels.items():
        if rel.reltype.endswith(("/notesSlide", "/tags")):
            continue
        if rel.is_external:
            rid_map[rid] = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            rid_map[rid] = dest.part.relate_to(rel.target_part, rel.reltype)

    for shape in source.shapes:
        new_el = copy.deepcopy(shape._element)
        _purger_refs_orphelines(new_el, rid_map)
        dest.shapes._spTree.append(new_el)
    return dest


def _reordonner(prs, slides_dans_lordre) -> None:
    """Réécrit l'ordre des slides ; celles absentes de la liste sont retirées.

    Une slide retirée l'est complètement : son entrée du sldIdLst ET sa relation
    depuis la présentation. Une relation orpheline vers une slide hors diaporama
    peut faire passer le fichier pour endommagé aux yeux de PowerPoint."""
    lst = prs.slides._sldIdLst
    elems = {int(s.get("id")): s for s in list(lst)}
    # slide.slide_id se résout via le sldIdLst : on capture les ids AVANT de le vider.
    ids_ordonnes = [s.slide_id for s in slides_dans_lordre]
    rids_retires = [e.get(qn("r:id")) for sid, e in elems.items() if sid not in ids_ordonnes]
    for s in list(lst):
        lst.remove(s)
    for sid in ids_ordonnes:
        lst.append(elems[sid])
    for rid in rids_retires:
        if rid and rid in prs.part.rels:
            prs.part.rels.pop(rid)


# --------------------------------------------------------------------------- #
#  Ouverture / contrôle du template
# --------------------------------------------------------------------------- #
def _ouvrir_gabarit() -> Presentation:
    if not template_disponible():
        raise FileNotFoundError(
            "Template PowerPoint introuvable (data/template_proposition.pptx). "
            "C'est un fichier local volontairement hors Git — voir le README."
        )
    prs = Presentation(config.TEMPLATE_PPTX_PATH)
    couverture = prs.slides[_IDX_COUVERTURE] if len(prs.slides) > _IDX_BUDGET else None
    texte_couverture = " ".join(
        s.text_frame.text for s in couverture.shapes if s.has_text_frame
    ) if couverture else ""
    if "Prestation Actuariat" not in texte_couverture:
        raise ValueError(
            "Le template ne correspond pas à la trame attendue (10 slides, page de "
            "garde « Prestation Actuariat »). Utilise le template Actuelia fourni."
        )
    return prs


# --------------------------------------------------------------------------- #
#  Remplissage des sections
# --------------------------------------------------------------------------- #
def _couverture(prs, demande: dict, redacteur: dict | None) -> None:
    slide = prs.slides[_IDX_COUVERTURE]
    reference = demande.get("reference") or ""
    titre = demande.get("titre") or ""
    ph18 = _placeholder(slide, 18)
    if ph18 is not None:
        _remplacer_marqueurs(ph18.text_frame, {
            "[Référence de l’appel d’offres]": reference,
            "[Intitulé de la mission]": titre,
        })
    ph19 = _placeholder(slide, 19)
    if ph19 is not None:
        _remplacer_marqueurs(ph19.text_frame, {"[JJ/MM/AAAA]": date.today().strftime("%d/%m/%Y")})

    ph20 = _placeholder(slide, 20)
    if ph20 is not None and redacteur:
        _remplacer_marqueurs(ph20.text_frame, {
            "[Prénom NOM]": redacteur.get("nom") or "[Prénom NOM]",
            "[Fonction]": redacteur.get("fonction") or "[Fonction]",
            "[prenom.nom]@actuelia.fr": redacteur.get("email") or "[prenom.nom]@actuelia.fr",
            "[06 XX XX XX XX]": redacteur.get("telephone") or "[06 XX XX XX XX]",
        })


def _contexte(prs, contenu: dict | None) -> None:
    contexte = (contenu or {}).get("contexte_redige", "").strip()
    if not contexte:
        return
    slide = prs.slides[_IDX_CONTEXTE]
    zone = _forme(slide, "Espace réservé du contenu 22")
    if zone is not None and zone.has_text_frame:
        _definir_texte(zone.text_frame, contexte)
        _ajuster_police_zone(zone)


def _retirer_frise(slide) -> None:
    """Retire la frise graphique des phases (ovales, pictos, groupes, zones de
    texte des phases) : elle ne tient pas des paragraphes de démarche."""
    for shape in list(slide.shapes):
        if shape.name in ("Titre 2", "Espace réservé du contenu 22") or shape.is_placeholder:
            continue
        shape._element.getparent().remove(shape._element)


def _modalites(prs, demande: dict, contenu: dict | None) -> None:
    slide = prs.slides[_IDX_MODALITES]
    intro = _forme(slide, "Espace réservé du contenu 22")
    if intro is None or not intro.has_text_frame:
        return
    _remplacer_marqueurs(intro.text_frame, {"[Client]": demande.get("client_nom") or "[Client]"})

    # La démarche est désormais isolée sur sa propre slide (_slide_demarche) pour
    # passer à l'échelle : la frise graphique du template débordait dès que les
    # descriptions de phase étaient longues. On la retire dès qu'une démarche est
    # rédigée (sinon on garde la frise et ses marqueurs [Phase N]).
    demarche = (contenu or {}).get("demarche", {}) or {}
    if any((demarche.get(p) or "").strip() for p in PHASES_DEMARCHE):
        _retirer_frise(slide)


def _titre_slide(slide):
    """Renvoie la forme titre d'une slide (par nom, sinon placeholder idx 0)."""
    forme = _forme(slide, "Titre 2")
    if forme is not None and forme.has_text_frame:
        return forme
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph
    return None


def _fixer_gras(run_el, gras: bool) -> None:
    rpr = run_el.find(qn("a:rPr"))
    if rpr is None:
        rpr = run_el.makeelement(qn("a:rPr"), {})
        run_el.insert(0, rpr)
    rpr.set("b", "1" if gras else "0")


def _para_phase(modele_p, label: str, desc: str):
    """Clone un paragraphe modèle en « label (gras) : desc (normal) »."""
    p = copy.deepcopy(modele_p)
    runs = p.findall(qn("a:r"))
    if not runs:
        t = p.find(".//" + qn("a:t"))
        if t is not None:
            t.text = f"{label}{desc}"
        return p
    for r in runs[1:]:
        r.getparent().remove(r)
    r_label = runs[0]
    r_desc = copy.deepcopy(r_label)
    r_label.find(qn("a:t")).text = label
    r_desc.find(qn("a:t")).text = desc
    _fixer_gras(r_label, True)
    _fixer_gras(r_desc, False)
    r_label.addnext(r_desc)
    return p


def _slide_demarche(prs, contenu: dict | None):
    """Crée une slide dédiée à la démarche opérationnelle (5 phases en texte).

    La frise graphique du template débordait dès que les descriptions de phase
    étaient longues et ne passait pas à l'échelle avec le nombre de consultants.
    On isole donc la démarche sur sa propre slide, calquée sur la charte de la
    slide Contexte (titre + zone de contenu). Renvoie la slide créée, ou None si
    la démarche n'a pas été rédigée.
    """
    demarche = (contenu or {}).get("demarche", {}) or {}
    phases = [(DEMARCHE_LABELS[p], (demarche.get(p) or "").strip()) for p in PHASES_DEMARCHE]
    phases = [(label, desc) for label, desc in phases if desc]
    if not phases:
        return None

    slide = _dupliquer_slide(prs, _IDX_CONTEXTE)

    titre = _titre_slide(slide)
    if titre is not None:
        _definir_texte(titre.text_frame, "Démarche opérationnelle")

    zone = _forme(slide, "Espace réservé du contenu 22")
    if zone is None or not zone.has_text_frame:
        return slide
    body = zone.text_frame._txBody
    paras = body.findall(qn("a:p"))
    if not paras:
        return slide
    modele = paras[0]
    for p in paras:
        body.remove(p)
    for label, desc in phases:
        body.append(_para_phase(modele, f"{label} : ", desc))
    _ajuster_police_zone(zone)
    return slide


def _charger_cv(ligne) -> dict:
    brut = _valeur(ligne, "cv_complet_json")
    if isinstance(brut, str) and brut:
        try:
            cv = json.loads(brut)
        except ValueError:
            cv = {}
    else:
        cv = brut or {}
    return cv if isinstance(cv, dict) else {}


def _textes_runs(p_element) -> list:
    return p_element.findall(".//" + qn("a:t"))


def _definir_runs(p_element, valeurs) -> None:
    """Fixe le texte des runs d'un paragraphe (None = run inchangé, garde son label)."""
    for t, valeur in zip(_textes_runs(p_element), valeurs):
        if valeur is not None:
            t.text = valeur


def _premieres_phrases(texte: str, nb: int = 2, max_caracteres: int = 340) -> str:
    """Les premières phrases d'un texte (pour un résumé court et borné)."""
    phrases = re.split(r"(?<=[.!?])\s+", (texte or "").strip())
    resume = " ".join(phrases[:nb]).strip()
    return resume if len(resume) <= max_caracteres else resume[:max_caracteres].rstrip() + "…"


def _equipe_projet(slide, lignes: list, *, retirer_pied_demarche: bool = False) -> None:
    """Remplit « L'équipe projet » : un bloc par consultant, bloc superviseur retiré.

    retirer_pied_demarche : retire le pied « Démarche opérationnelle proposée : »
    quand la démarche est portée par sa propre slide (sinon il reste orphelin).
    """
    if not lignes:
        return
    intro = _forme(slide, "Espace réservé du contenu 22")
    if intro is None or not intro.has_text_frame:
        return
    body = intro.text_frame._txBody
    paras = body.findall(qn("a:p"))

    def texte(p):
        return "".join((t.text or "") for t in p.findall(".//" + qn("a:t")))

    i_header = next((i for i, p in enumerate(paras) if "équipe projet" in texte(p)), None)
    i_footer = next((i for i, p in enumerate(paras) if "Démarche opérationnelle" in texte(p)), None)
    modele_id = next((p for p in paras if "Intervenant principal" in texte(p)), None)
    modele_exp = next((p for p in paras if texte(p).startswith("Expertise")), None)
    modele_va = next((p for p in paras if texte(p).startswith("Valeur ajoutée")), None)
    modele_role = next((p for p in paras if texte(p).startswith("Rôle :") and "[Responsab" in texte(p)), None)
    if i_header is None or i_footer is None or modele_id is None:
        return

    nouveaux = []
    for ligne in lignes:
        nom = f"{_valeur(ligne, 'prenom') or ''} {(_valeur(ligne, 'nom') or '').upper()}".strip()
        grade = _valeur(ligne, "grade") or _valeur(ligne, "seniorite") or ""
        xp = _valeur(ligne, "annees_experience")
        role = _valeur(ligne, "role_sur_mission")
        reste = f"– {grade}" if grade else "–"
        if xp:
            reste += f", {xp} ans d'expérience"
        if role:
            reste += f" ({role})"
        pid = copy.deepcopy(modele_id)
        _definir_runs(pid, [f"{nom} ", reste])
        nouveaux.append(pid)

        # Expertise = compétences (court) ; Valeur ajoutée = début de la synthèse
        # ciblée mission (borné à 2 phrases pour rester lisible quel que soit le
        # nombre de consultants) — la synthèse complète figure sur la fiche CV.
        competences = [c for c in _charger_cv(ligne).get("competences", []) if isinstance(c, str)]
        expertise = ", ".join(competences)
        if expertise and modele_exp is not None:
            pe = copy.deepcopy(modele_exp)
            _definir_runs(pe, [None, expertise])
            nouveaux.append(pe)

        valeur_ajoutee = _premieres_phrases(_valeur(ligne, "synthese_cv") or "")
        if valeur_ajoutee and modele_va is not None:
            pv = copy.deepcopy(modele_va)
            _definir_runs(pv, [None, valeur_ajoutee])
            nouveaux.append(pv)

        if role and modele_role is not None:
            pr = copy.deepcopy(modele_role)
            _definir_runs(pr, [None, role])
            nouveaux.append(pr)

    # Retire l'ancien bloc équipe (entre l'en-tête et le pied, bloc superviseur
    # compris) puis insère les nouveaux blocs, dans l'ordre.
    for p in paras[i_header + 1:i_footer]:
        body.remove(p)
    ref = paras[i_header]
    for p in nouveaux:
        ref.addnext(p)
        ref = p

    if retirer_pied_demarche and paras[i_footer].getparent() is not None:
        paras[i_footer].getparent().remove(paras[i_footer])
        # La frise a été retirée : la zone équipe peut occuper le bas de la slide.
        intro.height = Emu(int(9.1 * 914400))

    _ajuster_police_zone(intro)


def _fiche_cv(slide, ligne) -> None:
    def val(cle):
        try:
            return ligne[cle]
        except (KeyError, IndexError):
            return None

    prenom = val("prenom") or ""
    nom = (val("nom") or "").upper()
    titre = _forme(slide, "Titre 1")
    if titre is not None:
        _definir_texte(titre.text_frame, f"{prenom} {nom}".strip())

    grade = val("grade") or val("seniorite") or ""
    titre_poste = val("titre") or ""
    entete = ", ".join(x for x in (titre_poste, grade) if x)
    xp = f"{val('annees_experience')} années d'expérience" if val("annees_experience") else ""
    st = _placeholder(slide, 11)
    if st is not None:
        _definir_texte(st.text_frame, [entete or "[Titre, Grade]", xp])

    cv = _charger_cv(ligne)

    formation = val("formation") or cv.get("formation") or ""
    zt10 = _forme(slide, "ZoneTexte 10")
    if zt10 is not None and formation:
        _definir_texte(zt10.text_frame, formation)
        _reduire_si_deborde(zt10.text_frame)

    # Deux blocs d'expérience : bloc 1 (haut, ~top 4.4) et bloc 2 (bas, ~top 8.5).
    tables_shapes = sorted(_formes_nommees(slide, "Tableau 3"), key=lambda s: s.top or 0)
    groupes_exp = sorted(
        [g for g in slide.shapes if g.shape_type == 6 and any(
            "intitulé de la mi" in (s.text_frame.text if s.has_text_frame else "")
            for s in g.shapes)],
        key=lambda s: s.top or 0,
    )

    # Contenu des blocs : la synthèse ciblée mission (S3) en tête, puis les
    # expériences brutes du CV importé.
    blocs = []
    synthese = (val("synthese_cv") or "").strip()
    if synthese:
        blocs.append(("Synthèse pour la mission", synthese))
    for exp in cv.get("experiences", []):
        intitule = " — ".join(x for x in (exp.get("role"), exp.get("client")) if x) or "Expérience"
        blocs.append((intitule, exp.get("description") or ""))

    for i, (groupe, table_shape) in enumerate(zip(groupes_exp, tables_shapes)):
        if i < len(blocs):
            intitule, description = blocs[i]
            for sub in groupe.shapes:
                if sub.has_text_frame and "intitulé de la mi" in sub.text_frame.text:
                    _definir_texte(sub.text_frame, intitule)
            cellule = table_shape.table.cell(0, 0)
            _definir_texte(cellule.text_frame, description)
            _reduire_police_cellule(cellule.text_frame, len(description or ""))
        else:
            # Bloc sans contenu : on le retire pour ne pas laisser de marqueurs.
            groupe._element.getparent().remove(groupe._element)
            table_shape._element.getparent().remove(table_shape._element)

    # Photo et logos du collaborateur qui a servi d'exemple dans le template :
    # ils ne s'appliquent pas aux autres consultants -> on les retire.
    for shape in list(slide.shapes):
        if shape.shape_type == 13:  # PICTURE
            shape._element.getparent().remove(shape._element)


def _budget(prs, demande: dict, lignes: list) -> float:
    slide = prs.slides[_IDX_BUDGET]
    total = finance.total_mission(lignes)
    total_jours = sum(float(_valeur(l, "nb_jours") or 0) for l in lignes)
    intro = _forme(slide, "Espace réservé du contenu 22")
    if intro is not None and intro.has_text_frame:
        _remplacer_marqueurs(intro.text_frame, {
            "[Client]": demande.get("client_nom") or "[Client]",
            "[XX] jours": f"{total_jours:g} jours",
            "[XX XXX]": f"{total:,.0f}".replace(",", " "),
        })

    table_shape = next((s for s in slide.shapes if s.has_table), None)
    if table_shape is None:
        return total
    table = table_shape.table

    modele_tr = copy.deepcopy(table.rows[1]._tr)  # 2e ligne = ligne d'exemple
    # Retire toutes les lignes de données existantes (on garde l'entête).
    for row in list(table.rows)[1:]:
        row._tr.getparent().remove(row._tr)

    for i, ligne in enumerate(lignes, start=1):
        tr = copy.deepcopy(modele_tr)
        table._tbl.append(tr)
        cellules = table.rows[i].cells
        grade = ligne["grade"] if _a_valeur(ligne, "grade") else (ligne["seniorite"] if _a_valeur(ligne, "seniorite") else "")
        jours = ligne["nb_jours"] or 0
        tjm = ligne["tjm_applique"]
        _definir_texte(cellules[0].text_frame, str(grade or ""))
        _definir_texte(cellules[1].text_frame, f"{jours:g}")
        _definir_texte(cellules[2].text_frame,
                       f"{tjm:,.0f} €".replace(",", " ") if tjm else "—")

    return total


def _valeur(ligne, cle):
    try:
        return ligne[cle]
    except (KeyError, IndexError, TypeError):
        return None


def _a_valeur(ligne, cle) -> bool:
    return _valeur(ligne, cle) not in (None, "")


def _section_equipe(prs, lignes: list, *, slides_avant, slides_fin, demarche_slide=None) -> None:
    """Duplique la fiche CV pour chaque consultant, la remplit, et réordonne.

    Ordre final : intro (slides 0..6) + démarche (si rédigée) + fiches CV +
    budget/fin. La fiche CV modèle (index 7, avec ses marqueurs) est exclue dès
    qu'au moins un consultant est retenu. Les références de slides sont fournies
    par l'appelant, capturées AVANT toute duplication.
    """
    fiches = []
    for ligne in lignes:
        fiche = _dupliquer_slide(prs, _IDX_CV)
        _fiche_cv(fiche, ligne)
        fiches.append(fiche)

    milieu = [demarche_slide] if demarche_slide is not None else []
    if fiches:
        _reordonner(prs, slides_avant + milieu + fiches + slides_fin)
    elif milieu:
        # Démarche rédigée mais aucun consultant : on conserve la fiche CV modèle.
        _reordonner(prs, slides_avant + milieu + [prs.slides[_IDX_CV]] + slides_fin)


# --------------------------------------------------------------------------- #
#  Sauvegarde fidèle au template
# --------------------------------------------------------------------------- #
_CT_NS = "{http://schemas.openxmlformats.org/package/2006/content-types}"
_CT_SLIDE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"


def _canonique(donnees: bytes, nom: str):
    """Forme comparable d'une part : ignore la sérialisation (déclaration XML,
    entités, fins de ligne), pas le contenu. Les .rels sont comparées triées."""
    if not nom.endswith((".xml", ".rels")):
        return donnees
    try:
        arbre = etree.fromstring(donnees)
    except etree.XMLSyntaxError:
        return donnees
    if nom.endswith(".rels"):
        return b"|".join(sorted(etree.tostring(e) for e in arbre))
    return etree.tostring(arbre)


def _content_types_fideles(template_zip, noms_generes: set) -> bytes:
    """[Content_Types].xml du template, ajusté au minimum : retire les parts
    disparues, ajoute les slides créées. Tout le reste (Defaults svg compris)
    est conservé tel quel."""
    arbre = etree.fromstring(template_zip.read("[Content_Types].xml"))
    presentes = {"/" + n for n in noms_generes}
    couverts = set()
    for o in list(arbre.iter(_CT_NS + "Override")):
        if o.get("PartName") not in presentes:
            arbre.remove(o)
        else:
            couverts.add(o.get("PartName"))
    for n in sorted(noms_generes):
        pn = "/" + n
        if pn in couverts:
            continue
        # Les slides exigent leur Override (le Default « xml » ne suffit pas) ;
        # les autres nouveautés (.rels, médias) sont couvertes par les Defaults.
        if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
            o = etree.SubElement(arbre, _CT_NS + "Override")
            o.set("PartName", pn)
            o.set("ContentType", _CT_SLIDE)
    return etree.tostring(arbre, xml_declaration=True, encoding="UTF-8", standalone=True)


def _sauvegarder_fidele(prs, chemin_sortie) -> None:
    """Sauvegarde en préservant à l'octet près les parts non modifiées.

    python-pptx resérialise TOUT le paquet à la sauvegarde (déclarations XML,
    entités, ordre des entrées zip, content-types recomposés) ; PowerPoint
    s'est montré capricieux sur ce re-encodage, jusqu'à afficher blanches des
    slides du template pourtant jamais touchées. On ne réécrit donc que les
    parts réellement modifiées par la génération, tout le reste est recopié
    tel quel depuis le template, dans l'ordre d'origine."""
    tampon = io.BytesIO()
    prs.save(tampon)
    with zipfile.ZipFile(io.BytesIO(tampon.getvalue())) as genere, \
            zipfile.ZipFile(config.TEMPLATE_PPTX_PATH) as template, \
            zipfile.ZipFile(chemin_sortie, "w", zipfile.ZIP_DEFLATED) as sortie:
        noms_generes = set(genere.namelist())
        ecrits = set()
        for nom in template.namelist():
            if nom == "[Content_Types].xml":
                sortie.writestr(nom, _content_types_fideles(template, noms_generes))
            elif nom in noms_generes:
                brut_template = template.read(nom)
                brut_genere = genere.read(nom)
                identiques = _canonique(brut_template, nom) == _canonique(brut_genere, nom)
                sortie.writestr(nom, brut_template if identiques else brut_genere)
            else:
                continue  # part retirée (fiche CV modèle)
            ecrits.add(nom)
        for nom in genere.namelist():  # nouvelles parts (démarche, fiches CV)
            if nom not in ecrits and nom != "[Content_Types].xml":
                sortie.writestr(nom, genere.read(nom))


# --------------------------------------------------------------------------- #
#  Point d'entrée
# --------------------------------------------------------------------------- #
def generer_pptx(*, demande, lignes: list, chemin_sortie,
                 contenu: dict | None = None, redacteur: dict | None = None) -> float:
    """Génère le .pptx sur disque. Retourne le total mission (déterministe).

    demande / lignes acceptent des sqlite3.Row ou des dicts. contenu =
    contenu_genere_json (contexte + démarche S3). redacteur = coordonnées à
    porter en page de garde (nom, fonction, email, telephone).
    """
    demande = _en_dict(demande, ("titre", "reference", "client_nom"))
    prs = _ouvrir_gabarit()

    # Références des slides d'origine, capturées avant toute duplication.
    slides_avant = [prs.slides[i] for i in range(_IDX_CV)]  # 0..6 (Modalités inclus)
    slides_fin = [prs.slides[i] for i in range(_IDX_BUDGET, len(prs.slides))]  # 8..9

    _couverture(prs, demande, redacteur)
    _contexte(prs, contenu)
    _modalites(prs, demande, contenu)
    demarche_slide = _slide_demarche(prs, contenu)
    _equipe_projet(prs.slides[_IDX_MODALITES], lignes,
                   retirer_pied_demarche=demarche_slide is not None)
    total = _budget(prs, demande, lignes)
    _section_equipe(prs, lignes, slides_avant=slides_avant, slides_fin=slides_fin,
                    demarche_slide=demarche_slide)

    _sauvegarder_fidele(prs, chemin_sortie)
    return total


def _en_dict(source, cles) -> dict:
    if isinstance(source, dict):
        return source
    resultat = {}
    for cle in cles:
        try:
            resultat[cle] = source[cle]
        except (KeyError, IndexError, TypeError):
            resultat[cle] = None
    return resultat
